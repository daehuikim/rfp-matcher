"""정량 '사람이 쓸법한' 평가 하네스 — 정답파일 vs 추출 결과.

지표 4종(문서별):
 1) 탭수      : 정답 시트수(빈시트 제외) vs 추출 distinct 탭수
 2) 행수      : 정답 데이터행 vs 추출 행수 (ratio 1.0 이 human-like)
 3) 상세 글자수: median / p90 / max 나란히
 4) ID 커버리지: 총괄표 선언 건수(DECLARED) vs 추출 distinct 문서ID (누락/여분 목록)
추가 검증:
 5) 페이지 단조성: source_page 역행 횟수 (목표 0)
 6) 분할행 동일ID: 같은 code 를 공유하는 행 통계 (재번호 금지 확인)

사용:
  backend/.venv/bin/python3 temp/human_metrics.py            # temp/overnight/*_reqs.json 전부
  backend/.venv/bin/python3 temp/human_metrics.py 양형 기아  # 이름 부분일치 필터
입력: temp/overnight/<stem>_reqs.json (배치 스크립트가 저장한 API /requirements 응답)
"""
from __future__ import annotations

import json
import re
import statistics
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
OVERNIGHT = ROOT / "temp" / "overnight"
ANSWERS = ROOT / "data_real" / "raw" / "answers"

# 문서 자체 요구사항 총괄표의 선언 건수 (원문 실측)
DECLARED = {
    "양형": {"ECR": 9, "COM": 2, "SFR": 47, "CSR": 2, "PER": 4,
             "INR": 8, "DAR": 15, "TER": 7, "SER": 15, "QUR": 9},
    "법제처": {"SFR": 21, "DAR": 26, "PER": 10, "SIR": 10, "TER": 6,
               "SER": 18, "QUR": 6, "COR": 14, "PMR": 26, "PSR": 14},
}

# 정답파일 매핑 (stem 부분일치)
GOLD_FILES = {
    "기아": ANSWERS / "(보안) 기아 차세대 고객센터 구축 사업_조견표.xlsx",
    "대한항공": ANSWERS / "대한항공 조견표_0130 _2240 (1).pptx",
    "법제처": ANSWERS / "법제처_요구사항_정리_개선.xlsx",
}

_DOC_ID = re.compile(r"^([A-Z]{2,4})(?:[-–][A-Z]{2,6})?(?:[-–]\d{1,4}|\d{2,4})$")


def _dist(lengths: list[int]) -> str:
    if not lengths:
        return "-"
    ls = sorted(lengths)
    p90 = ls[min(len(ls) - 1, int(len(ls) * 0.9))]
    return f"{int(statistics.median(ls))}/{p90}/{max(ls)}"


def gold_stats(name: str) -> dict | None:
    p = None
    for k, v in GOLD_FILES.items():
        if k in name:
            p = v
    if p is None or not p.exists():
        return None
    if p.suffix == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(p, read_only=True, data_only=True)
        tabs = 0
        rows = 0
        lengths: list[int] = []
        for ws in wb.worksheets:
            grid = [[("" if c is None else str(c)) for c in r] for r in ws.iter_rows(values_only=True)]
            grid = [r for r in grid if any(x.strip() for x in r)]
            if len(grid) < 2:
                continue
            # 총괄표 시트(ID부여규칙 XXX-OOO, 문자 O 플레이스홀더)는 데이터 아님.
            # 숫자 0 을 클래스에 넣으면 SFR-001 같은 진짜 ID 도 걸려 전 시트가
            # 제외되는 함정(기아 1573→590 실측) — 문자 O/Ø 만, 뒤에 숫자 없어야 함.
            flat = " ".join(x for r in grid[:12] for x in r)
            if re.search(r"[A-Z]{2,4}\s*[-–]\s*[OØ]{2,3}(?![0-9])", flat) or "총괄" in ws.title:
                continue
            tabs += 1
            # 상세열 = 평균 길이 최장 열
            ncol = max(len(r) for r in grid)
            avg = [0.0] * ncol
            for c in range(ncol):
                vals = [len(r[c].strip()) for r in grid[1:] if c < len(r) and r[c].strip()]
                avg[c] = sum(vals) / len(vals) if vals else 0
            det = max(range(ncol), key=lambda c: avg[c])
            for r in grid[1:]:
                v = r[det].strip() if det < len(r) else ""
                if v:
                    rows += 1
                    lengths.append(len(v))
        return {"tabs": tabs, "rows": rows, "len": _dist(lengths)}
    if p.suffix == ".pptx":
        from pptx import Presentation
        pres = Presentation(p)
        rows = 0
        lengths: list[int] = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if not shape.has_table:
                    continue
                for r in shape.table.rows:
                    cells = [c.text.strip() for c in r.cells]
                    longest = max((len(c) for c in cells), default=0)
                    if longest >= 10:
                        rows += 1
                        lengths.append(longest)
        return {"tabs": None, "rows": rows, "len": _dist(lengths)}
    return None


def eval_doc(reqs_path: Path) -> dict:
    name = reqs_path.stem.replace("_reqs", "")
    data = json.loads(reqs_path.read_text(encoding="utf-8"))
    rows = [x["requirement"] for x in data]
    cats: list[str] = []
    for r in rows:
        c = r.get("category") or "?"
        if c not in cats:
            cats.append(c)
    lengths = [len((r.get("detail") or "").strip()) for r in rows if (r.get("detail") or "").strip()]
    # 페이지 단조성 — 탭 내부 기준(탭 우선 정렬이라 탭 경계의 페이지 점프는 정상)
    backward = 0
    prev_by_tab: dict = {}
    for r in rows:
        p = r.get("source_page")
        if p is None:
            continue
        t = r.get("category") or ""
        if t in prev_by_tab and p < prev_by_tab[t]:
            backward += 1
        prev_by_tab[t] = p
    pages = [r.get("source_page") for r in rows]
    # ID 커버리지
    decl = None
    for k, v in DECLARED.items():
        if k in name:
            decl = v
    idcov = None
    if decl:
        by_pfx: dict[str, set[str]] = {}
        for r in rows:
            code = (r.get("code") or "").strip()
            m = _DOC_ID.match(code)
            if m and m.group(1) in decl:
                by_pfx.setdefault(m.group(1), set()).add(code)
        idcov = {}
        for pfx, want in decl.items():
            got = len(by_pfx.get(pfx, set()))
            idcov[pfx] = (got, want)
    # 분할행 동일ID (문서ID 행만)
    from collections import Counter
    doc_codes = [r.get("code") for r in rows if _DOC_ID.match((r.get("code") or "").strip())]
    cc = Counter(doc_codes)
    multi = sum(1 for v in cc.values() if v > 1)
    return {
        "name": name, "rows": len(rows), "tabs": len(cats),
        "len": _dist(lengths), "backward": backward,
        "page_none": sum(1 for p in pages if p is None),
        "doc_id_rows": len(doc_codes), "doc_id_distinct": len(cc), "doc_id_multi": multi,
        "idcov": idcov, "tabs_list": cats,
    }


def main() -> None:
    pats = sys.argv[1:]
    files = sorted(OVERNIGHT.glob("*_reqs.json"))
    if pats:
        files = [f for f in files if any(p in f.stem for p in pats)]
    if not files:
        print(f"입력 없음: {OVERNIGHT}/*_reqs.json")
        return
    out_lines: list[str] = []

    def emit(s: str = "") -> None:
        print(s)
        out_lines.append(s)

    emit("| 문서 | 추출행(정답행) | 추출탭(정답탭) | 상세 med/p90/max (정답) | 페이지역행 | 문서ID 행/고유/반복 |")
    emit("|---|---|---|---|---|---|")
    details: list[str] = []
    for f in files:
        m = eval_doc(f)
        g = gold_stats(m["name"]) or {}
        emit(f"| {m['name']} | {m['rows']}({g.get('rows','-')}) | {m['tabs']}({g.get('tabs','-')}) "
             f"| {m['len']} ({g.get('len','-')}) | {m['backward']} "
             f"| {m['doc_id_rows']}/{m['doc_id_distinct']}/{m['doc_id_multi']} |")
        if m["idcov"]:
            total_got = sum(g_ for g_, _ in m["idcov"].values())
            total_want = sum(w for _, w in m["idcov"].values())
            cov = " ".join(f"{p}:{g_}/{w}" for p, (g_, w) in m["idcov"].items())
            details.append(f"**{m['name']} ID커버리지 {total_got}/{total_want}** — {cov}")
        details.append(f"{m['name']} 탭: {', '.join(m['tabs_list'][:20])}"
                       + (" …" if len(m["tabs_list"]) > 20 else ""))
    emit()
    for d in details:
        emit(d)
    (OVERNIGHT / "metrics.md").write_text("\n".join(out_lines), encoding="utf-8")
    print(f"\n→ {OVERNIGHT / 'metrics.md'}")


if __name__ == "__main__":
    main()
