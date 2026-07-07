"""정답지(gold) 기반 추출 자동 평가 — recall(빠진 요구)·precision(junk) 실측.

gold 요구사항 텍스트 vs run_v_rule_reqs 추출 상세를 퍼지매칭(문자 3-gram Jaccard + 포함).
- recall   = 매칭된 gold / 전체 gold   (요구사항 얼마나 잡았나)
- precision = 매칭된 추출 / 전체 추출   (추출 중 진짜 요구 비율; 낮으면 junk 많음)
junk 샘플(gold에 없는 추출)·missed 샘플(추출에 없는 gold)을 보여줘 정조준 튜닝.

사용: backend/.venv/bin/python temp/dr_eval.py [기아|대한항공 ...]
gold 추가: GOLD dict 에 항목만 넣으면 다른 파일도 동일 평가.
"""
import sys, re, os
from pathlib import Path
ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "backend"))
os.chdir(ROOT / "backend")

RAW = Path("../data_real/raw")
ANS = RAW / "answers"
GOLD = {
    "기아": ("pdf_기아자동차.pdf", "xlsx", ANS / "(보안) 기아 차세대 고객센터 구축 사업_조견표.xlsx"),
    "대한항공": ("pdf_대한항공.pdf", "pptx", ANS / "대한항공 조견표_0130 _2240 (1).pptx"),
    # src 가 ../ 로 시작하면 RAW 기준이 아니라 backend cwd 기준 경로로 해석
    "법제처": ("../data/raw/법제처_생성형 AI 법령검색 시스템 구축_RFI_20260116.hwpx", "xlsx",
             ANS / "법제처_요구사항_정리_개선.xlsx"),
}

_MARK = re.compile(r"^[\s]*(?:[❍○●◦▪▫◆◇■□▣◈▶▷▸‣⁃∙·•⦁*※\-–—]|\d+[.)]|[가-힣][.)]|\([0-9가-힣]+\))+")
def norm(s: str) -> str:
    s = _MARK.sub("", str(s or "").strip())
    return re.sub(r"[^가-힣a-z0-9]", "", s.lower())

def grams(s: str) -> set:
    return {s[i:i+3] for i in range(len(s) - 2)} or ({s} if s else set())

def load_gold_xlsx(path: Path) -> list[str]:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    out = []
    id_rule_pat = re.compile(r"\b[A-Z]{2,4}[-–][O0Ø]{2,3}\b")   # 'SFR-OOO' 총괄표 ID부여규칙 신호
    for sn in wb.sheetnames:
        ws = wb[sn]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        flat = " ".join(str(c) for r in rows for c in r if c)
        if len(id_rule_pat.findall(flat)) >= 3:
            continue   # 요구사항 총괄표(집계/범례 시트) — 실제 요구사항 행이 아니므로 제외
        # 요건/요구사항 열 찾기(헤더). 여러 열이 매칭되면(예: '요구사항 분류' vs '요구사항
        # 상세설명-세부내용') **평균 셀 길이가 가장 긴 열**을 채택 — 분류/라벨열은 짧고
        # 반복적(시트당 거의 동일값)이라 그걸 고르면 중복제거 후 gold 가 왕창 사라진다
        # (법제처 실측: '요구사항 분류'열 선택 시 151개→14개로 붕괴).
        hdr = [str(c or "") for c in rows[0]]
        ncol = max((len(r) for r in rows), default=0)

        def _avglen(c: int) -> float:
            vals = [str(r[c]) for r in rows[1:] if c < len(r) and r[c]]
            return sum(len(v) for v in vals) / len(vals) if vals else 0.0

        # '상세/세부내용' 같은 구체적 신호를 '요구사항/요건'(분류열에도 흔함)보다 우선 후보로.
        specific = [i for i, h in enumerate(hdr) if re.search(r"상세|세부\s*내용", h)]
        generic = [i for i, h in enumerate(hdr) if re.search(r"요건|요구사항|요구 사항|내용", h)]
        candidates = specific or generic
        col = max(candidates, key=_avglen) if candidates else None
        body = rows[1:] if col is not None else rows
        if col is None:
            if not ncol:
                continue
            col = max(range(ncol), key=_avglen)
        for r in body:
            if col < len(r) and r[col] and len(str(r[col]).strip()) >= 6:
                out.append(str(r[col]).strip())
    return out

def load_gold_pptx(path: Path) -> list[str]:
    from pptx import Presentation
    prs = Presentation(str(path))
    out = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if not shp.has_table:
                continue
            tbl = shp.table
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                # 요구사항 텍스트 = 그 행에서 가장 긴 셀(≥10자)
                long = [c for c in cells if len(c) >= 10]
                if long:
                    out.append(max(long, key=len))
    return out

def evaluate(key: str):
    src_name, kind, gold_path = GOLD[key]
    if not gold_path.exists():
        print(f"[{key}] gold 없음: {gold_path}"); return
    gold_raw = load_gold_xlsx(gold_path) if kind == "xlsx" else load_gold_pptx(gold_path)
    # gold 정규화·중복제거(짧은 것 제외)
    gold = {}
    for g in gold_raw:
        n = norm(g)
        if len(n) >= 6:
            gold[n] = g
    from prototype.v_rule.adapter import run_v_rule_reqs
    src = Path(src_name) if src_name.startswith("../") else RAW / src_name
    reqs = run_v_rule_reqs(src, Path("/tmp/dreval") / key)
    ext = []
    for r in reqs:
        d = getattr(r, "detail", "") or ""
        n = norm(d)
        if len(n) >= 4:
            ext.append((n, d))
    gset = {n: grams(n) for n in gold}
    eset = [(n, grams(n), d) for n, d in ext]

    def best(a_grams, a_norm, pool):
        bs = 0.0
        for n2, g2 in pool:
            if a_norm in n2 or n2 in a_norm:
                return 1.0
            inter = len(a_grams & g2)
            if inter:
                j = inter / len(a_grams | g2)
                if j > bs:
                    bs = j
        return bs

    TH = 0.5
    gpool = [(n, g) for n, g in gset.items()]
    epool = [(n, g) for n, g, _ in eset]
    gold_hit = {n: best(g, n, epool) >= TH for n, g in gset.items()}
    matched_gold = sum(gold_hit.values())
    matched_ext = [(best(g, n, gpool) >= TH, d) for n, g, d in eset]
    n_ext_ok = sum(1 for ok, _ in matched_ext if ok)
    recall = matched_gold / max(1, len(gold))
    precision = n_ext_ok / max(1, len(eset))
    print(f"\n{'='*66}\n[{key}] gold {len(gold)}개 / 추출 {len(eset)}개")
    print(f"  recall(요구 잡음)   = {recall*100:5.1f}%  ({matched_gold}/{len(gold)})")
    print(f"  precision(junk아님) = {precision*100:5.1f}%  ({n_ext_ok}/{len(eset)})  → junk {100-precision*100:.1f}%")
    junk = [d for ok, d in matched_ext if not ok]
    print(f"  --- JUNK 추출 샘플(gold에 없음) {len(junk)}개 중 15 ---")
    for d in junk[:15]:
        print(f"     ✗ {d[:72]}")
    print("  --- (참고) MISSED gold 샘플 10 ---")
    miss = [gold[n] for n, hit in gold_hit.items() if not hit]
    for g in miss[:10]:
        print(f"     · {g[:72]}")

if __name__ == "__main__":
    keys = [k for k in sys.argv[1:] if k in GOLD] or list(GOLD)
    for k in keys:
        try:
            evaluate(k)
        except Exception as e:
            import traceback; print(f"[{k}] 실패: {e}"); traceback.print_exc()
